import re
import io
from pptx import Presentation
from langchain_text_splitters import TokenTextSplitter
import tiktoken

def pptx_process_to_chunks(
    file_bytes: bytes, file_hash: str, chunk_size: int = 512, chunk_overlap: int = 64
) -> list[dict]:
    """
    Extracts text from PPTX slides and splits into chunks.
    Each chunk is mapped to a slide number.
    """
    prs = Presentation(io.BytesIO(file_bytes))

    splitter = TokenTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        encoding_name="cl100k_base",
    )

    final_chunks = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        slide_text_elements = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                slide_text_elements.append(text)

        if not slide_text_elements:
            continue

        clean_text = " ".join(slide_text_elements)
        clean_text = re.sub(r"\s+", " ", clean_text).strip()

        if not clean_text:
            continue

        # Split and store
        chunks = splitter.split_text(clean_text)
        for chunk in chunks:
            final_chunks.append(
                {
                    "text": chunk,
                    "file_hash": file_hash,
                    "page_number": slide_num, # Using page_number to maintain compatibility with DB schema
                }
            )

    return final_chunks

