import tiktoken
import os
import hashlib
import io
from uuid import uuid4
from datetime import datetime, timezone

from fastapi import HTTPException, UploadFile
import fitz
from pptx import Presentation
from app.services.embedding_service import EmbeddingService
from app.services.vector_db_service import VectorDBService
from app.tasks import process_document

from supabase import Client

BUCKET = "pdfs"  # Define the storage bucket name

# Supported formats mapping: MIME type -> extension
SUPPORTED_FORMATS = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
}


class UploadService:
    def __init__(
        self,
        vector_service: VectorDBService,
        db_client: Client,
        embedding_service: EmbeddingService,
    ):
        self.vector_service = vector_service
        self.db = db_client
        self.embedding_service = embedding_service

    async def execute(
        self,
        file: UploadFile,
        user_id: str,
    ):
        # 1) Determine file extension from MIME type or filename
        file_ext = SUPPORTED_FORMATS.get(file.content_type)
        if not file_ext:
            # Fallback to filename extension (useful for generic MIME types like octet-stream)
            name_ext = os.path.splitext(file.filename or "")[1].lower()
            if name_ext in SUPPORTED_FORMATS.values():
                file_ext = name_ext
            else:
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported file type: {file.content_type}. Please upload PDF or PPTX.",
                )

        contents = await file.read()
        if not contents:
            raise HTTPException(status_code=400, detail="Empty file")

        # Check file size limit (50MB)
        MAX_FILE_SIZE = 50 * 1024 * 1024
        file_size = len(contents)
        if file_size > MAX_FILE_SIZE:
            raise HTTPException(
                status_code=413,
                detail=f"File size exceeds 50MB limit. File size: {len(contents) / (1024 * 1024):.2f}MB",
            )

        page_count = 0
        full_text = ""

        try:
            if file_ext == ".pdf":
                with fitz.open(stream=contents, filetype="pdf") as doc:
                    page_count = len(doc)
                    for page_num in range(page_count):
                        page = doc.load_page(page_num)
                        full_text += page.get_text("text") + "\n"
            elif file_ext == ".pptx":
                prs = Presentation(io.BytesIO(contents))
                page_count = len(prs.slides)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            full_text += shape.text + " "
                    full_text += "\n"
        except Exception as e:
            raise HTTPException(
                status_code=400, detail=f"Failed to parse {file_ext} content: {str(e)}"
            )

        encoding = tiktoken.get_encoding("cl100k_base")
        token_count = len(encoding.encode(full_text))

        # Check limits
        MAX_PAGES = int(os.getenv("MAX_PDF_PAGES", "2000"))
        if page_count > MAX_PAGES:
            raise HTTPException(
                status_code=413,
                detail=f"File exceeds the {MAX_PAGES}-page limit. It has {page_count} pages.",
            )

        MAX_DOCUMENT_TOKENS = int(os.getenv("MAX_DOCUMENT_TOKENS", "1000000"))
        if token_count > MAX_DOCUMENT_TOKENS:
            raise HTTPException(
                status_code=413,
                detail=f"File exceeds the {MAX_DOCUMENT_TOKENS}-token limit. It contains {token_count} tokens.",
            )

        sha256 = hashlib.sha256(contents).hexdigest()
        file_name = file.filename or f"upload{file_ext}"

        # 1. Lookup: Check for an existing document by 'file_hash'
        existing = (
            self.db.table("documents")
            .select("*")
            .eq("file_hash", sha256)
            .maybe_single()
            .execute()
        )

        document = None
        should_trigger_task = False
        message = ""
        restored_from_trash = False

        # Case 1: Document exists and is NOT failed (ready, processing, indexing)
        if existing and existing.data and existing.data.get("status") != "failed":
            document = existing.data
            document_id = document["id"]
            message = "Document already exists."
        else:
            # Case 2: Document is new OR previously failed
            document_id = existing.data["id"] if (existing and existing.data) else str(uuid4())
            storage_path = f"{document_id}{file_ext}"

            doc_data = {
                "id": document_id,
                "file_hash": sha256,
                "file_name": file_name,
                "file_path": storage_path,
                "file_size": file_size,
                "page_count": page_count,
                "token_count": token_count,
                "status": "pending",
                "created_at": datetime.now(timezone.utc).isoformat(),
            }

            # Idempotent upsert on the document record
            upsert_res = self.db.table("documents").upsert(doc_data, on_conflict="file_hash").execute()
            if not upsert_res.data:
                raise HTTPException(status_code=500, detail="Failed to upsert document row")
            
            document = upsert_res.data[0]

            try:
                # Get correct MIME type for storage upload
                content_type = file.content_type
                if content_type not in SUPPORTED_FORMATS:
                    content_type = next(k for k, v in SUPPORTED_FORMATS.items() if v == file_ext)

                # Idempotent storage upload (x-upsert: true)
                res = self.db.storage.from_(BUCKET).upload(
                    path=storage_path,
                    file=contents,
                    file_options={
                        "content-type": content_type,
                        "x-upsert": "true",
                    },
                )

                if hasattr(res, "error") and res.error is not None:
                    raise Exception(f"Supabase Storage Error: {res.error}")

                should_trigger_task = True
                message = "Upload successful." if not existing else "Recovered failed document and re-started processing."

            except Exception as e:
                # If storage fails, mark it as failed so it can be recovered later
                self.db.table("documents").update({"status": "failed"}).eq("id", document_id).execute()
                raise HTTPException(status_code=500, detail=f"Storage upload failed: {e}")

        # 4. User Linking:
        # - If an existing link is trashed, restore it (deleted_at -> null)
        # - If no link exists, insert a fresh link
        link_res = (
            self.db.table("user_library")
            .select("deleted_at")
            .eq("user_id", user_id)
            .eq("document_id", document_id)
            .maybe_single()
            .execute()
        )

        if link_res and link_res.data:
            if link_res.data.get("deleted_at") is not None:
                self.db.table("user_library").update({"deleted_at": None}).eq(
                    "user_id", user_id
                ).eq("document_id", document_id).execute()
                restored_from_trash = True
        else:
            self.db.table("user_library").insert(
                {"user_id": user_id, "document_id": document_id}
            ).execute()

        # 5. Single Trigger: if 'should_trigger_task' is True, await the task exactly once
        if should_trigger_task:
            await process_document(
                document_id=document_id,
                file_hash=sha256,
                file_path=document["file_path"],
                db_client=self.db,
                vector_service=self.vector_service,
                embedding_service=self.embedding_service,
            )

            # Re-fetch document to get the latest status (likely 'ready' or 'failed')
            final_doc_res = self.db.table("documents").select("*").eq("id", document_id).maybe_single().execute()
            if final_doc_res.data:
                document = final_doc_res.data

        if restored_from_trash and not should_trigger_task:
            message = "Document restored from trash."

        return {
            "document": {
                "id": document["id"],
                "name": document["file_name"],
                "status": document["status"],
                "size": document["file_size"],
                "pageCount": document["page_count"],
                "uploadedAt": document.get("created_at"),
            },
            "message": message,
        }
